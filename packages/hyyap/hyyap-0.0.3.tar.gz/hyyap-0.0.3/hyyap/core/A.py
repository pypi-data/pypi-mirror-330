from docx import Document
from docx.shared import Pt, Inches, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_ALIGN_VERTICAL


class WordDocx:
    def __init__(self):
        self.doc = Document()

    def title(self, title_text, level=1, font_size=20):
        """向文档中添加标题"""
        title = self.doc.add_heading(title_text, level=level)
        title.runs[0].font.size = Pt(font_size)
        return self

    def paragraph(self, paragraph_text, font_size=12):
        """向文档中添加段落"""
        paragraph = self.doc.add_paragraph(paragraph_text)
        paragraph.runs[0].font.size = Pt(font_size)
        return self

    def bullet_list(self, items):
        """向文档中添加无序列表"""
        for item in items:
            self.doc.add_paragraph(item, style='List Bullet')
        return self

    def numbered_list(self, items):
        """向文档中添加有序列表"""
        for item in items:
            self.doc.add_paragraph(item, style='List Number')
        return self

    def picture(self, picture_path, width=4):
        """向文档中添加图片"""
        self.doc.add_picture(picture_path, width=Inches(width))
        return self

    def table(self, rows, cols, headers, data):
        """向文档中添加表格"""
        table = self.doc.add_table(rows=rows, cols=cols)
        hdr_cells = table.rows[0].cells
        for i, header in enumerate(headers):
            hdr_cells[i].text = header
        for row_num, row_data in enumerate(data):
            row_cells = table.rows[row_num + 1].cells
            for col_num, cell_data in enumerate(row_data):
                row_cells[col_num].text = cell_data
        return self

    def table_column_widths(self, table_index, column_widths):
        """设置表格指定列的宽度"""
        if table_index < len(self.doc.tables):
            table = self.doc.tables[table_index]
            for col_index, width in enumerate(column_widths):
                if col_index < len(table.columns):
                    table.columns[col_index].width = Cm(width)
        return self

    def table_row_height(self, table_index, row_index, line_spacing=2.0):
        """设置表格指定行的行高（通过段落行距调整）"""
        if table_index < len(self.doc.tables):
            table = self.doc.tables[table_index]
            if row_index < len(table.rows):
                for cell in table.rows[row_index].cells:
                    paragraph = cell.paragraphs[0]
                    paragraph_format = paragraph.paragraph_format
                    paragraph_format.line_spacing = line_spacing
                    cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        return self

    def save(self, file_path):
        """保存文档，确保保存为 .docx 格式"""
        if not file_path.lower().endswith('.docx'):
            if '.' in file_path:
                file_path = file_path.rsplit('.', 1)[0] + '.docx'
            else:
                file_path = file_path + '.docx'
        self.doc.save(file_path)

from openpyxl import Workbook
from openpyxl.styles import Font
from openpyxl.utils import get_column_letter


class ExcelXlsx:
    def __init__(self):
        # 初始化一个新的 Excel 工作簿
        self.wb = Workbook()
        # 获取默认的工作表
        self.ws = self.wb.active

    def add_title(self, title_text, row=1, col=1):
        """向 Excel 工作表中添加标题"""
        cell = self.ws.cell(row=row, column=col, value=title_text)
        cell.font = Font(size=20, bold=True)
        return self

    def add_paragraph(self, paragraph_text, row, col):
        """向 Excel 工作表中添加段落文本"""
        cell = self.ws.cell(row=row, column=col, value=paragraph_text)
        cell.font = Font(size=12)
        return self

    def add_list(self, items, start_row, col, is_bullet=True):
        """向 Excel 工作表中添加列表"""
        for i, item in enumerate(items, start=start_row):
            cell = self.ws.cell(row=i, column=col)
            if is_bullet:
                cell.value = f"• {item}"
            else:
                cell.value = f"{i - start_row + 1}. {item}"
            cell.font = Font(size=12)
        return self

    def add_picture(self):
        # 注意：openpyxl 对图片操作相对复杂，这里暂不实现，可考虑使用其他库如 pandas 结合 xlsxwriter
        print("Excel 中添加图片暂不支持此简单实现。")
        return self

    def add_table(self, headers, data, start_row, start_col):
        """向 Excel 工作表中添加表格"""
        # 写入表头
        for col_num, header in enumerate(headers, start=start_col):
            cell = self.ws.cell(row=start_row, column=col_num)
            cell.value = header
            cell.font = Font(size=12, bold=True)
        # 写入表格数据
        for row_num, row_data in enumerate(data, start=start_row + 1):
            for col_num, cell_data in enumerate(row_data, start=start_col):
                cell = self.ws.cell(row=row_num, column=col_num)
                cell.value = cell_data
        return self

    def set_table_column_widths(self, start_col, column_widths):
        """设置表格指定列的宽度"""
        for col_index, width in enumerate(column_widths, start=start_col):
            col_letter = get_column_letter(col_index)
            self.ws.column_dimensions[col_letter].width = width
        return self

    def set_table_row_height(self, start_row, row_indexes, heights):
        """设置表格指定行的行高"""
        for i, row_index in enumerate(row_indexes, start=start_row):
            self.ws.row_dimensions[row_index].height = heights[i]
        return self

    def save(self, file_path):
        """保存 Excel 文档，确保保存为 .xlsx 格式"""
        if not file_path.lower().endswith('.xlsx'):
            if '.' in file_path:
                file_path = file_path.rsplit('.', 1)[0] + '.xlsx'
            else:
                file_path = file_path + '.xlsx'
        self.wb.save(file_path)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


class PPTGenerator:
    def __init__(self):
        # 初始化一个新的 PPT 演示文稿
        self.prs = Presentation()

    def add_title_slide(self, title_text, subtitle_text=""):
        """添加标题幻灯片"""
        title_slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = title_text
        subtitle.text = subtitle_text
        return self

    def add_content_slide(self, title_text, content_items):
        """添加带有内容列表的幻灯片"""
        content_slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        body_shape = slide.shapes.placeholders[1]
        title.text = title_text
        tf = body_shape.text_frame
        for item in content_items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
        return self

    def add_table_slide(self, title_text, headers, data):
        """添加带有表格的幻灯片"""
        table_slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(table_slide_layout)
        title = slide.shapes.title
        title.text = title_text
        rows, cols = len(data) + 1, len(headers)
        left = top = Inches(2)
        width = Inches(6)
        height = Inches(0.8)
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        # 写入表头
        for col_index, header in enumerate(headers):
            table.cell(0, col_index).text = header
            table.cell(0, col_index).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        # 写入表格数据
        for row_index, row_data in enumerate(data, start=1):
            for col_index, cell_data in enumerate(row_data):
                table.cell(row_index, col_index).text = cell_data
        return self

    def add_picture_slide(self, title_text, picture_path):
        """添加带有图片的幻灯片"""
        picture_slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(picture_slide_layout)
        title = slide.shapes.title
        title.text = title_text
        left = top = Inches(2)
        slide.shapes.add_picture(picture_path, left, top, width=Inches(6))
        return self

    def save(self, file_path):
        """保存 PPT 文档，确保保存为 .pptx 格式"""
        if not file_path.lower().endswith('.pptx'):
            if '.' in file_path:
                file_path = file_path.rsplit('.', 1)[0] + '.pptx'
            else:
                file_path = file_path + '.pptx'
        self.prs.save(file_path)